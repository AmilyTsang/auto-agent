import jinja2
# import weasyprint
import pandas as pd
import json
from datetime import datetime
import os
from typing import Dict, Any, List

class ReportGenerator:
    """报告生成类"""
    
    def __init__(self, template_dir: str = "templates"):
        self.template_dir = template_dir
        # 确保报告目录存在
        os.makedirs("reports", exist_ok=True)
        # 确保模板目录存在
        os.makedirs(template_dir, exist_ok=True)
        
        # 创建Jinja2环境
        self.template_env = jinja2.Environment(
            loader=jinja2.FileSystemLoader(template_dir),
            autoescape=jinja2.select_autoescape(['html', 'xml'])
        )
    
    def generate_html_report(self, data: Dict[str, Any], title: str = "数据分析报告") -> str:
        """生成HTML报告"""
        # 生成报告ID
        report_id = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        report_path = f"reports/{report_id}.html"
        
        # 准备模板数据
        template_data = {
            'title': title,
            'generated_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'data': data
        }
        
        # 渲染HTML模板
        html_content = self._render_template('report_template.html', template_data)
        
        # 保存报告
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return report_path
    
    def generate_pdf_report(self, data: Dict[str, Any], title: str = "数据分析报告") -> str:
        """生成PDF报告"""
        # 暂时返回HTML报告路径，因为WeasyPrint需要系统依赖
        return self.generate_html_report(data, title)
    
    def generate_pptx_report(self, data: Dict[str, Any], title: str = "数据分析报告") -> str:
        """生成PowerPoint报告"""
        # 这里使用python-pptx库
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return "错误：需要安装python-pptx库"
        
        # 创建演示文稿
        prs = Presentation()
        
        # 添加标题幻灯片
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # 添加数据概览幻灯片
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "数据概览"
        
        # 添加分析结果幻灯片
        if 'descriptive_analysis' in data:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = "描述性分析"
        
        # 保存演示文稿
        pptx_id = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        pptx_path = f"reports/{pptx_id}.pptx"
        prs.save(pptx_path)
        
        return pptx_path
    
    def _render_template(self, template_name: str, context: Dict[str, Any]) -> str:
        """渲染模板"""
        # 检查模板是否存在
        template_path = os.path.join(self.template_dir, template_name)
        if not os.path.exists(template_path):
            # 创建默认模板
            self._create_default_template(template_name)
        
        # 渲染模板
        template = self.template_env.get_template(template_name)
        return template.render(**context)
    
    def _create_default_template(self, template_name: str):
        """创建默认模板"""
        if template_name == 'report_template.html':
            template_content = """
<!DOCTYPE html>
<html>
<head>
    <title>{{ title }}</title>
    <style>
        body {
            font-family: Arial, sans-serif;
            margin: 20px;
            line-height: 1.6;
        }
        h1 {
            color: #333;
            text-align: center;
        }
        h2 {
            color: #555;
            border-bottom: 2px solid #eee;
            padding-bottom: 10px;
        }
        h3 {
            color: #666;
        }
        table {
            border-collapse: collapse;
            width: 100%;
            margin: 20px 0;
        }
        th, td {
            border: 1px solid #ddd;
            padding: 8px;
            text-align: left;
        }
        th {
            background-color: #f2f2f2;
            font-weight: bold;
        }
        tr:nth-child(even) {
            background-color: #f9f9f9;
        }
        .summary {
            background-color: #f5f5f5;
            padding: 15px;
            border-radius: 5px;
            margin: 20px 0;
        }
        .chart {
            margin: 20px 0;
            text-align: center;
        }
        .generated-at {
            text-align: right;
            color: #999;
            font-size: 0.9em;
            margin-top: 40px;
        }
    </style>
</head>
<body>
    <h1>{{ title }}</h1>
    <div class="generated-at">生成时间：{{ generated_at }}</div>
    
    {% if data.get('data_overview') %}
    <h2>数据概览</h2>
    <div class="summary">
        <p>数据形状：{{ data.data_overview.shape }}</p>
        <p>列数：{{ data.data_overview.columns|length }}</p>
        <p>行数：{{ data.data_overview.rows }}</p>
    </div>
    {% endif %}
    
    {% if data.get('descriptive_analysis') %}
    <h2>描述性分析</h2>
    <h3>基本统计信息</h3>
    <table>
        <tr>
            <th>列名</th>
            <th>计数</th>
            <th>均值</th>
            <th>标准差</th>
            <th>最小值</th>
            <th>最大值</th>
        </tr>
        {% for col, stats in data.descriptive_analysis.basic_stats.items() %}
        <tr>
            <td>{{ col }}</td>
            <td>{{ stats.count }}</td>
            <td>{{ "%.2f"|format(stats.mean) if stats.mean is not none else "N/A" }}</td>
            <td>{{ "%.2f"|format(stats.std) if stats.std is not none else "N/A" }}</td>
            <td>{{ "%.2f"|format(stats.min) if stats.min is not none else "N/A" }}</td>
            <td>{{ "%.2f"|format(stats.max) if stats.max is not none else "N/A" }}</td>
        </tr>
        {% endfor %}
    </table>
    
    <h3>缺失值分析</h3>
    <table>
        <tr>
            <th>列名</th>
            <th>缺失值数量</th>
            <th>缺失值比例</th>
        </tr>
        {% for col, count in data.descriptive_analysis.missing_values.items() %}
        <tr>
            <td>{{ col }}</td>
            <td>{{ count }}</td>
            <td>{{ "%.2f%%"|format(count / data.data_overview.rows * 100) }}</td>
        </tr>
        {% endfor %}
    </table>
    {% endif %}
    
    {% if data.get('visualizations') %}
    <h2>数据可视化</h2>
    {% for viz in data.visualizations %}
    <div class="chart">
        <h3>{{ viz.title }}</h3>
        <img src="{{ viz.path }}" alt="{{ viz.title }}" style="max-width: 100%;">
    </div>
    {% endfor %}
    {% endif %}
    
    {% if data.get('insights') %}
    <h2>分析洞察</h2>
    <div class="summary">
        {% for insight in data.insights %}
        <p>{{ insight }}</p>
        {% endfor %}
    </div>
    {% endif %}
</body>
</html>
            """
            with open(os.path.join(self.template_dir, template_name), 'w', encoding='utf-8') as f:
                f.write(template_content)
